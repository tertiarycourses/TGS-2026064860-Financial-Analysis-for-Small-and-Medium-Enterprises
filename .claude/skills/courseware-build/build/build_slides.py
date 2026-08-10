#!/usr/bin/env python3
"""Generate the WSQ Financial Analysis for SMEs course slide deck
(all-white Tertiary house style, highly visual).

Content is copied from the legacy v10 master deck and re-set in the house
design system (tile grids, flow diagrams, cards, profile cards, images).
Hands-on activities appear as overview slides only — the detailed step-by-step
lives in the Learner Guide.
"""
import os, sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

HERE = os.path.dirname(os.path.abspath(__file__))
sys.path.insert(0, HERE)
import course_data as C
from data_domain1 import DOMAIN1
from data_domain2 import DOMAIN2
from data_domain3 import DOMAIN3
ACTIVITIES = DOMAIN1 + DOMAIN2 + DOMAIN3

def _find_repo(start):
    env = os.environ.get("COURSE_REPO")
    if env and os.path.isdir(env):
        return env
    d = start
    for _ in range(8):
        d = os.path.dirname(d)
        if os.path.isdir(os.path.join(d, "courseware")) and os.path.isdir(os.path.join(d, "activities")):
            return d
    return os.path.dirname(os.path.dirname(HERE))
REPO = _find_repo(HERE)
ASSETS = os.path.join(os.path.dirname(HERE), "assets")        # skill assets (logo)
CW_ASSETS = os.path.join(REPO, "courseware", "assets")        # course content images

# ---------------- palette ----------------
BLUE=RGBColor(0x1F,0x6F,0xEB); TEAL=RGBColor(0x10,0xB9,0x81); AMBER=RGBColor(0xF5,0x9E,0x0B)
INK=RGBColor(0x16,0x1B,0x26); GREY=RGBColor(0x5B,0x63,0x72); LIGHT=RGBColor(0xF5,0xF8,0xFC)
WHITE=RGBColor(0xFF,0xFF,0xFF); LINE=RGBColor(0xE2,0xE8,0xF0); VIOLET=RGBColor(0x7C,0x3A,0xED)

prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
SW,SH=prs.slide_width,prs.slide_height
BLANK=prs.slide_layouts[6]

def slide(): return prs.slides.add_slide(BLANK)
def rect(s,x,y,w,h,color,line=None):
    sp=s.shapes.add_shape(1,x,y,w,h); sp.fill.solid(); sp.fill.fore_color.rgb=color
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb=line; sp.line.width=Pt(1)
    sp.shadow.inherit=False; return sp
def oval(s,x,y,w,h,color):
    sp=s.shapes.add_shape(9,x,y,w,h); sp.fill.solid(); sp.fill.fore_color.rgb=color
    sp.line.fill.background(); sp.shadow.inherit=False; return sp
def txt(s,x,y,w,h,runs,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP,space=4):
    tb=s.shapes.add_textbox(x,y,w,h); tf=tb.text_frame; tf.word_wrap=True; tf.vertical_anchor=anchor
    for i,line in enumerate(runs):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.alignment=align; p.space_after=Pt(space)
        for t,sz,col,bold in line:
            r=p.add_run(); r.text=t; r.font.size=Pt(sz); r.font.bold=bold
            r.font.color.rgb=col; r.font.name="Arial"
    return tb
def bullets(s,x,y,w,h,items,size=18,color=INK,gap=10,mcolor=BLUE):
    tb=s.shapes.add_textbox(x,y,w,h); tf=tb.text_frame; tf.word_wrap=True
    for i,it in enumerate(items):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.space_after=Pt(gap)
        lvl=it[1] if isinstance(it,tuple) else 0
        text=it[0] if isinstance(it,tuple) else it
        r=p.add_run(); r.text=("•  " if lvl==0 else "–  ")+text
        r.font.size=Pt(size if lvl==0 else size-2); r.font.color.rgb=color if lvl==0 else GREY
        r.font.name="Arial"; r.font.bold=(lvl==0 and isinstance(it,tuple) and len(it)>2 and it[2])
    return tb

PAGE={"n":0}
def footer(s):
    PAGE["n"]+=1
    txt(s,Inches(0.4),Inches(7.05),Inches(7.5),Inches(0.35),
        [[(f"{C.SHORT_TITLE}  ·  {C.COURSE_CODE}",9,GREY,False)]])
    txt(s,Inches(5.0),Inches(7.05),Inches(3.3),Inches(0.35),
        [[("© 2026 Tertiary Infotech Academy Pte Ltd",9,GREY,False)]],align=PP_ALIGN.CENTER)
    txt(s,Inches(12.4),Inches(7.05),Inches(0.6),Inches(0.35),
        [[(str(PAGE["n"]),9,GREY,False)]],align=PP_ALIGN.RIGHT)
def head(s,title,kicker=None,kcolor=BLUE):
    rect(s,0,0,SW,SH,WHITE); rect(s,0,0,Inches(0.28),Inches(1.55),kcolor)
    if kicker: txt(s,Inches(0.85),Inches(0.5),Inches(11.6),Inches(0.4),[[(kicker,14,kcolor,True)]])
    txt(s,Inches(0.85),Inches(0.9),Inches(11.9),Inches(0.9),[[(title,29,INK,True)]])
    rect(s,Inches(0.85),Inches(1.7),Inches(11.63),Inches(0.02),LINE)
    return s
def _logo(name):
    p=os.path.join(ASSETS,name)
    return p if os.path.exists(p) else None
def _content_img(name):
    p=os.path.join(CW_ASSETS,name)
    return p if os.path.exists(p) else None

# ---------------- slide templates ----------------
def cover():
    s=slide(); rect(s,0,0,SW,SH,WHITE)
    rect(s,0,0,SW,Inches(0.22),BLUE); rect(s,0,Inches(7.28),SW,Inches(0.22),TEAL)
    org=_logo("tertiary-infotech-logo.png")
    if org: s.shapes.add_picture(org,Inches(0.85),Inches(0.7),height=Inches(1.05))
    rect(s,Inches(10.6),Inches(0.72),Inches(1.95),Inches(1.0),BLUE)
    txt(s,Inches(10.6),Inches(0.86),Inches(1.95),Inches(0.5),[[("WSQ",22,WHITE,True)]],align=PP_ALIGN.CENTER)
    txt(s,Inches(10.6),Inches(1.32),Inches(1.95),Inches(0.4),[[("FINANCIAL ANALYSIS",9,WHITE,True)]],align=PP_ALIGN.CENTER)
    txt(s,Inches(0.9),Inches(2.3),Inches(12),Inches(0.6),[[("COURSE SLIDES  ·  WSQ",16,BLUE,True)]])
    txt(s,Inches(0.9),Inches(2.85),Inches(12.0),Inches(1.9),[[(C.TITLE,40,INK,True)]])
    rect(s,Inches(0.92),Inches(4.75),Inches(2.4),Inches(0.06),TEAL)
    txt(s,Inches(0.9),Inches(5.05),Inches(12),Inches(1.4),
        [[(f"WSQ Course Code: {C.COURSE_CODE}",16,GREY,False)],
         [("Conducted by Tertiary Infotech Academy Pte Ltd  ·  UEN 201200696W",14,GREY,False)]],space=6)
    txt(s,Inches(0.9),Inches(6.5),Inches(12),Inches(0.4),[[(f"Version {C.VERSION}  ·  {C.VERSION_DATE}",12,GREY,False)]])
    txt(s,Inches(0.9),Inches(6.85),Inches(12),Inches(0.34),[[("© 2026 Tertiary Infotech Academy Pte Ltd. All rights reserved.  ·  www.tertiarycourses.com.sg",10,GREY,False)]])

def section(kicker,title,n,sub=""):
    s=slide(); rect(s,0,0,SW,SH,WHITE); rect(s,0,0,Inches(0.28),SH,BLUE)
    rect(s,Inches(0.85),Inches(2.5),Inches(0.14),Inches(2.0),TEAL)
    txt(s,Inches(1.25),Inches(2.55),Inches(11),Inches(0.6),[[(kicker,18,BLUE,True)]])
    txt(s,Inches(1.25),Inches(3.0),Inches(11.4),Inches(1.6),[[(title,40,INK,True)]])
    if sub: txt(s,Inches(1.27),Inches(4.55),Inches(11),Inches(0.8),[[(sub,16,GREY,False)]])
    txt(s,Inches(10.0),Inches(0.7),Inches(2.8),Inches(1.6),[[(n,72,RGBColor(0xE2,0xE8,0xF0),True)]],align=PP_ALIGN.RIGHT)
    footer(s)
def content(title,items,kicker=None,size=20):
    s=head(slide(),title,kicker); bullets(s,Inches(0.85),Inches(1.95),Inches(11.6),Inches(4.9),items,size=size); footer(s); return s
def two_col(title,left,right,kicker=None,lhead="",rhead=""):
    s=head(slide(),title,kicker)
    rect(s,Inches(0.85),Inches(1.95),Inches(5.7),Inches(4.7),LIGHT); rect(s,Inches(6.95),Inches(1.95),Inches(5.55),Inches(4.7),LIGHT)
    if lhead: txt(s,Inches(1.1),Inches(2.15),Inches(5.2),Inches(0.4),[[(lhead,16,BLUE,True)]])
    if rhead: txt(s,Inches(7.2),Inches(2.15),Inches(5.0),Inches(0.4),[[(rhead,16,TEAL,True)]])
    bullets(s,Inches(1.1),Inches(2.7),Inches(5.2),Inches(3.8),left,size=16)
    bullets(s,Inches(7.2),Inches(2.7),Inches(5.05),Inches(3.8),right,size=16,mcolor=TEAL); footer(s); return s
def cards3(title,cards,kicker):
    s=head(slide(),title,kicker); xs=[Inches(0.85),Inches(5.0),Inches(9.15)]
    for i,c in enumerate(cards[:3]):
        x=xs[i]; col=c[0]
        rect(s,x,Inches(1.95),Inches(3.65),Inches(4.7),LIGHT); rect(s,x,Inches(1.95),Inches(3.65),Inches(0.12),col)
        txt(s,x+Inches(0.25),Inches(2.2),Inches(3.2),Inches(0.6),[[(c[1],19,col,True)]])
        bullets(s,x+Inches(0.25),Inches(2.95),Inches(3.2),Inches(3.4),c[2],size=14,mcolor=col,gap=9)
    footer(s); return s
def big_statement(line1,line2,kicker,color=BLUE):
    s=slide(); rect(s,0,0,SW,SH,WHITE); rect(s,0,0,Inches(0.28),SH,color)
    txt(s,Inches(1.1),Inches(2.2),Inches(11),Inches(0.5),[[(kicker,16,color,True)]])
    txt(s,Inches(1.1),Inches(2.8),Inches(11.3),Inches(2.4),[[(line1,38,INK,True)]])
    if line2: txt(s,Inches(1.12),Inches(4.9),Inches(11),Inches(1.2),[[(line2,20,GREY,False)]])
    footer(s); return s
import math
PALETTE=[BLUE,TEAL,VIOLET,AMBER]
def tile_grid(title,items,kicker=None,cols=2,size=15,icons=None,accent=BLUE):
    s=head(slide(),title,kicker,kcolor=accent)
    n=len(items); rows=math.ceil(n/cols)
    X0=Inches(0.85); Y0=Inches(1.95); TOTW=Inches(11.63); AREAH=Inches(4.78)
    gx=Inches(0.3); gy=Inches(0.26)
    cw=int((TOTW-gx*(cols-1))/cols); ch=int((AREAH-gy*(rows-1))/rows)
    bd=Inches(0.6)
    for i,it in enumerate(items):
        r=i//cols; c=i%cols
        x=int(X0+(cw+gx)*c); y=int(Y0+(ch+gy)*r); col=PALETTE[i%len(PALETTE)]
        rect(s,x,y,cw,ch,LIGHT); rect(s,x,y,Inches(0.1),ch,col)
        oval(s,x+Inches(0.28),int(y+ch/2-bd/2),bd,bd,col)
        ic=icons[i] if icons else str(i+1)
        txt(s,x+Inches(0.28),int(y+ch/2-bd/2),bd,bd,[[(ic,19,WHITE,True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
        tx=x+Inches(1.08); tw=cw-Inches(1.32)
        if isinstance(it,tuple):
            txt(s,tx,int(y+Inches(0.14)),tw,int(ch-Inches(0.2)),
                [[(it[0],size+2,INK,True)],[(it[1],size-2,GREY,False)]],anchor=MSO_ANCHOR.MIDDLE,space=3)
        else:
            txt(s,tx,int(y+Inches(0.1)),tw,int(ch-Inches(0.16)),[[(it,size,INK,False)]],anchor=MSO_ANCHOR.MIDDLE)
    footer(s); return s
def flow_h(title,steps,kicker=None,color=BLUE):
    s=head(slide(),title,kicker,kcolor=color)
    n=len(steps); X0=Inches(0.85); TOTW=Inches(11.63); gap=Inches(0.34)
    cw=int((TOTW-gap*(n-1))/n); y=Inches(2.55); ch=Inches(3.15); bd=Inches(0.82)
    for i,st in enumerate(steps):
        x=int(X0+(cw+gap)*i)
        rect(s,x,y,cw,ch,LIGHT); rect(s,x,y,cw,Inches(0.1),color)
        oval(s,int(x+cw/2-bd/2),int(y+Inches(0.42)),bd,bd,color)
        txt(s,int(x+cw/2-bd/2),int(y+Inches(0.42)),bd,bd,[[(str(i+1),30,WHITE,True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
        txt(s,x+Inches(0.16),int(y+Inches(1.55)),cw-Inches(0.32),int(ch-Inches(1.7)),[[(st,14,INK,False)]],align=PP_ALIGN.CENTER)
        if i<n-1:
            txt(s,int(x+cw-Inches(0.04)),int(y+ch/2-Inches(0.3)),int(gap+Inches(0.08)),Inches(0.6),
                [[("▶",15,color,True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    footer(s); return s
def trainer_slide(kicker,name,role,rows,initials,accent=BLUE):
    s=head(slide(),"About the Trainer",kicker,kcolor=accent)
    lx=Inches(0.85); lw=Inches(3.65)
    rect(s,lx,Inches(1.95),lw,Inches(4.7),LIGHT); rect(s,lx,Inches(1.95),lw,Inches(0.12),accent)
    bd=Inches(1.7); ax=int(lx+(lw-bd)/2)
    oval(s,ax,Inches(2.5),bd,bd,accent)
    txt(s,ax,Inches(2.5),bd,bd,[[(initials,44,WHITE,True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    txt(s,lx+Inches(0.15),Inches(4.55),lw-Inches(0.3),Inches(0.6),[[(name,21,INK,True)]],align=PP_ALIGN.CENTER)
    txt(s,lx+Inches(0.15),Inches(5.2),lw-Inches(0.3),Inches(1.2),[[(role,13,GREY,False)]],align=PP_ALIGN.CENTER)
    rx=Inches(4.9); rw=Inches(7.6); ry=Inches(1.95); rh=Inches(4.7)
    n=len(rows); gy=Inches(0.2); th=int((rh-gy*(n-1))/n)
    for i,(label,val) in enumerate(rows):
        y=int(ry+(th+gy)*i); col=PALETTE[i%len(PALETTE)]
        rect(s,rx,y,rw,th,LIGHT); rect(s,rx,y,Inches(0.1),th,col)
        vruns=[(val,14,INK,False)] if val else [("____________________________________________",13,LINE,False)]
        txt(s,rx+Inches(0.32),y,rw-Inches(0.6),th,
            [[(label.upper(),11,col,True)],vruns],anchor=MSO_ANCHOR.MIDDLE,space=3)
    footer(s); return s
def activity_overview(tag,title,desc,build,services,kicker):
    s=head(slide(),title,kicker,kcolor=TEAL)
    rect(s,Inches(0.85),Inches(1.85),Inches(2.2),Inches(0.5),TEAL)
    txt(s,Inches(0.85),Inches(1.9),Inches(2.2),Inches(0.4),[[(tag,16,WHITE,True)]],align=PP_ALIGN.CENTER)
    txt(s,Inches(0.85),Inches(2.55),Inches(11.7),Inches(1.6),[[(desc,19,INK,False)]])
    rect(s,Inches(0.85),Inches(4.35),Inches(11.7),Inches(1.95),LIGHT)
    txt(s,Inches(1.1),Inches(4.55),Inches(11),Inches(0.4),[[("You'll produce",14,BLUE,True)]])
    txt(s,Inches(1.1),Inches(4.95),Inches(11.2),Inches(0.8),[[(build,17,INK,True)]])
    txt(s,Inches(1.1),Inches(5.75),Inches(11.2),Inches(0.5),[[("Tools:  ",13,GREY,True),(services,13,GREY,False)],
        [("Detailed step-by-step: see this activity in the Learner Guide.",12,TEAL,True)]],space=3)
    footer(s); return s
def img_slide(title,img,kicker=None,caption="",h=Inches(4.75),accent=BLUE):
    """Full-width content image under the standard header."""
    s=head(slide(),title,kicker,kcolor=accent)
    p=_content_img(img)
    if p:
        from PIL import Image as _PIL
        iw,ih=_PIL.open(p).size
        maxw=Inches(11.6); maxh=h
        scale=min(maxw/iw, maxh/ih)
        w=int(iw*scale); hh=int(ih*scale)
        x=int((SW-w)/2); y=int(Inches(1.95)+(maxh-hh)/2)
        s.shapes.add_picture(p,x,y,width=w,height=hh)
    if caption:
        txt(s,Inches(0.85),Inches(6.72),Inches(11.6),Inches(0.35),[[(caption,11,GREY,False)]],align=PP_ALIGN.CENTER)
    footer(s); return s
def img_text_slide(title,img,items,kicker=None,accent=BLUE,size=15):
    """Left bullets, right image."""
    s=head(slide(),title,kicker,kcolor=accent)
    bullets(s,Inches(0.85),Inches(2.05),Inches(5.6),Inches(4.6),items,size=size)
    p=_content_img(img)
    if p:
        from PIL import Image as _PIL
        iw,ih=_PIL.open(p).size
        maxw=Inches(6.0); maxh=Inches(4.6)
        scale=min(maxw/iw, maxh/ih)
        w=int(iw*scale); hh=int(ih*scale)
        x=int(Inches(6.7)+(maxw-w)/2); y=int(Inches(2.0)+(maxh-hh)/2)
        s.shapes.add_picture(p,x,y,width=w,height=hh)
    footer(s); return s
def formula_slide(title,formulas,kicker=None,accent=BLUE,notes=None):
    """formulas: list of (name, formula, note). Rendered as wide formula cards."""
    s=head(slide(),title,kicker,kcolor=accent)
    n=len(formulas); Y0=Inches(2.0); AREAH=Inches(4.65); gy=Inches(0.24)
    ch=int((AREAH-gy*(n-1))/n)
    for i,(name,formula,note) in enumerate(formulas):
        y=int(Y0+(ch+gy)*i); col=PALETTE[i%len(PALETTE)]
        rect(s,Inches(0.85),y,Inches(11.63),ch,LIGHT); rect(s,Inches(0.85),y,Inches(0.1),ch,col)
        txt(s,Inches(1.15),y,Inches(3.6),ch,[[(name,15,col,True)]],anchor=MSO_ANCHOR.MIDDLE)
        txt(s,Inches(4.9),y,Inches(4.4),ch,[[(formula,15,INK,True)]],anchor=MSO_ANCHOR.MIDDLE)
        txt(s,Inches(9.45),y,Inches(2.9),ch,[[(note,11.5,GREY,False)]],anchor=MSO_ANCHOR.MIDDLE)
    footer(s); return s
def brk(kind,dur,color=AMBER):
    s=slide(); rect(s,0,0,SW,SH,WHITE)
    rect(s,0,0,SW,Inches(0.22),color); rect(s,0,Inches(7.28),SW,Inches(0.22),color)
    rect(s,Inches(5.4),Inches(2.35),Inches(2.53),Inches(0.1),color)
    txt(s,0,Inches(2.75),SW,Inches(1.2),[[(kind,48,INK,True)]],align=PP_ALIGN.CENTER)
    txt(s,0,Inches(4.05),SW,Inches(0.8),[[(dur,22,color,True)]],align=PP_ALIGN.CENTER); PAGE["n"]+=1

# ============================================================ BUILD
MARKS={}
def mark(key):
    """Record the page number of the NEXT slide to be added (cover = slide 1)."""
    MARKS[key]=len(prs.slides._sldIdLst)+1

mark("cover")
cover()

# ---------------- ADMIN ----------------
mark("admin")
section("COURSE ADMINISTRATION","Welcome & Housekeeping","")
content("Digital Attendance (Mandatory)",[
 "It is mandatory to take the AM, PM and Assessment digital attendance for WSQ-funded courses.",
 "The trainer or administrator will show you the digital attendance QR code generated from the SSG portal.",
 "Scan the QR code with your mobile phone camera and submit your attendance.",
 "A minimum of 75% attendance is required to be eligible for assessment and funding."],kicker="TRAQOM · SSG DIGITAL ATTENDANCE")
trainer_slide("YOUR TRAINER · GENERAL","Your Trainer","General Trainer template —\nto be completed by the trainer",
 [("Name",""),("Title / Designation",""),("Qualifications",""),
  ("Areas of expertise",""),("Training & industry experience",""),("Contact","")],
 initials="?",accent=GREY)
trainer_slide("YOUR TRAINER",C.TRAINER,"Principal Trainer\nTertiary Infotech Academy Pte Ltd",
 [("Role","Principal Trainer, Tertiary Infotech Academy Pte Ltd"),
  ("Background","PhD; ACTA/ACLP-certified trainer with deep experience in data analytics, finance and business applications."),
  ("Delivers","WSQ courses on financial analysis, accounting systems, data analytics and business software."),
  ("Founder","Founder and lead instructor at Tertiary Infotech / Tertiary Courses.")],
 initials="AA",accent=BLUE)
content("Let's Know Each Other",[
 "Your name and organisation / role.",
 "Your experience with accounting or finance (if any).",
 "What you want to be able to analyse or decide after this course."],kicker="ICE-BREAKER")
tile_grid("Ground Rules",[
 "Set your mobile phone to silent mode.","Participate actively — no question is stupid.",
 "Mutual respect: agree to disagree.","One conversation at one time.",
 "Be punctual; return from breaks on time.","75% attendance is required."],
 kicker="HOUSEKEEPING",cols=2,size=15)
img_text_slide("Download Course Material",
 "lms-login.png",
 [("Go to https://lms-tms.tertiaryinfotech.com",0),
  ("Log in with your registered email — an OTP is sent to you.",0),
  ("Open your course: "+C.TITLE,0),
  ("Download the slides, Learner Guide and activity worksheets.",0),
  ("Keep them handy — the final assessment is open book.",0)],
 kicker="LMS / TMS COURSE PORTAL",size=16)
two_col("Lesson Plan — 2 Days, 8 Hours/Day",[
 (f"Day 1 — {C.DAY_THEMES[1]}",0),
 ("Digital Attendance (AM)",1),
 ("Trainer and learner introductions · Learning outcomes · Course outline",1),
 ("Topic 1: Understanding Financial Statements (Activity 1)",1),
 ("Lunch break · Digital Attendance (PM)",1),
 ("Topic 2: Analysing Financial Ratios (Activities 2–3)",1)],
 [(f"Day 2 — {C.DAY_THEMES[2]}",0),
 ("Digital Attendance (AM)",1),
 ("Topic 3: Planning & Budgeting using Financial Statements (Activities 4–6)",1),
 ("Lunch break · Digital Attendance (PM)",1),
 ("Revision · Course feedback & TRAQOM survey",1),
 ("Digital Attendance (Assessment) · Final Assessment (WA + PP)",1)],
 kicker="SCHEDULE",lhead="Day 1",rhead="Day 2")
tile_grid("Skills Framework (TSC)",[
 ("TSC Title · Code",f"{C.TSC_TITLE}  ·  {C.TSC_CODE}"),
 ("Knowledge K1–K2","Statement of financial position · Balance sheet"),
 ("Knowledge K3–K4","Income and cash flow statements · Statement of changes in equity"),
 ("Knowledge K5","Financial statement analysis techniques"),
 ("Ability A1","Identify trends by comparing ratios across time periods and statement types"),
 ("Ability A2","Prepare and interpret performance and position using financial statements")],
 kicker="WSQ SKILLS FRAMEWORK",cols=2,size=13)
tile_grid("Learning Outcomes",[
 ("LO1 — Understand","Financial statements: the balance sheet, income statement and cash flow statement."),
 ("LO2 — Evaluate","The organisation's financial performance from the trend of financial ratios."),
 ("LO3 — Analyze","Financial statements and prepare the organisation's position.")],
 kicker="WHAT YOU'LL ACHIEVE",cols=1,size=16)
cards3("Course Outline — 3 Topics",[
 (BLUE,"Topic 1 (K1–K3)",["Overview of Finance and Chart of Accounts","Balance Sheet Statement","Profit and Loss (P&L) Statement","Cash Flow Statement"]),
 (TEAL,"Topic 2 (K4, A1)",["Ratios for Corporate Profitability","Ratios for Corporate Performance","Equity Changes Statement"]),
 (VIOLET,"Topic 3 (K5, A2)",["Analyse Financial Statements","Financial Planning","Capital Budgeting"])],kicker="COURSE OUTLINE")
content("Criteria for Funding",[
 "Minimum attendance rate of 75% based on the SSG Digital Attendance record.",
 "Complete the assessment and be assessed as 'Competent'."],kicker="WSQ FUNDING")
content("Briefing for Assessment",[
 "Place phones and other materials under the table or on the floor.",
 "No photos or recording of assessment scripts.","No discussion during the assessment.",
 "Use a black/blue pen for hard-copy assessments.","No liquid paper or correction tape.",
 "Assessment scripts are collected when the time is up."])
content("Assessment",[
 C.ASSESSMENT["written"], C.ASSESSMENT["practical"],
 "Format: Open Book — slides, Learner Guide and approved materials only.",
 C.ASSESSMENT["note"],"An appeal process is available if required."],kicker="FINAL ASSESSMENT")
flow_h("Assessment Flow",[
 "TRAQOM survey — scan the QR code on the LMS",
 "Assessment digital attendance — scan the SSG QR",
 "Sit WA (SAQ) then PP — open book",
 "Submit your answers on the LMS",
 "Sign the Assessment Summary Record"],kicker="ON ASSESSMENT DAY")

# ================================================== TOPIC 1
T=C.TOPICS[0]
mark("topic1")
section(f"TOPIC {T['code']}",T["title"],T["code"],T["subtitle"])
tile_grid(f"Key Concepts — {T['title']}",T["concepts"],kicker=f"TSC {T['weighting']}",cols=2,size=14)

big_statement("Finance is attaining money, investing it, and gaining a return on investment.",
 "Any business decision has financial impact — finance is the common language of the boardroom.","BUSINESS FINANCE")
tile_grid("Significance of Finance",[
 "Any business decision has financial impact.","It is the common language used in board meetings.",
 "The more you understand finance, the more the insights.","Dealing with finance people becomes easy.",
 "Relate finance to your own business unit.","Drive and achieve better performance."],
 kicker="WHY FINANCE MATTERS",cols=2,size=14)
tile_grid("The Finance Business Unit",[
 ("FP&A","Budgeting, forecasting, actual vs variance, management reports."),
 ("Operations","Purchases and payables · expenses · sales and receivables."),
 ("Controls","Internal controls, audit, processes and policies."),
 ("Systems","Finance tools, transformation, ERP implementation, data & BI."),
 ("M&A","Due diligence, checklists, M&A budget, pre/post-integration."),
 ("Tax & Treasury","Corporate tax, GST, tax audit and cash flow management.")],
 kicker="CORPORATE ORG STRUCTURE · FINANCE",cols=2,size=13)
two_col("The Financial System",[
 ("Regulators & institutions",0),
 ("Ministry of Finance · ACRA · IRAS · SAC · CPF",1),
 ("Banks · insurance companies · mutual funds · private equity funds",1),
 ("Lenders provide capital to the system",1)],
 [("Borrowers & businesses",0),
 ("Sole proprietorship · partnership / LLP",1),
 ("Private company · public listed company",1),
 ("Non-profit organisations · individuals · foreign companies",1)],
 kicker="LENDERS ↔ BORROWERS",lhead="Lenders & regulators",rhead="Borrowers")
tile_grid("Accounting — the Big Picture",[
 ("Three questions","How much money do I need? How do I raise it? How do I manage it?"),
 ("The equation","Assets = Liabilities + Owner's Equity — always in balance."),
 ("The mechanics","Double entry · transaction analysis · summary of transactions."),
 ("The outputs","Income statement · owner's equity statement · balance sheet · statement of cash flows.")],
 kicker="ACCOUNTING",cols=2,size=14)
two_col("Objectives & Advantages of Accounting",[
 ("Objectives",0),
 ("Maintain a systematic record of transactions",1),
 ("Ascertain profit or loss and financial position",1),
 ("Assist management and communicate information to users",1),
 ("Prevent manipulation and fraud",1)],
 [("Advantages",0),
 ("Financial information about the business; replaces memory",1),
 ("Facilitates benchmarking, tax settlement and loans",1),
 ("Evidence in court; facilitates sale of the business",1),
 ("Helps in decision making",1)],
 kicker="WHY WE ACCOUNT",lhead="Objectives",rhead="Advantages")
tile_grid("Types of Accounting",[
 ("Financial accounting","Aggregates financial information into external reports — IFRS, GAAP."),
 ("Management accounting","Internal operational reporting — cost accounting, payroll, payables."),
 ("Tax accounting","Compliance with tax regulations, filings and future tax planning."),
 ("Auditing","Assurance that financial statements fairly present results and position."),
 ("Forensic accounting","Reconstruction of financial information when records are incomplete."),
 ("Internal auditing","Examination of systems and transactions to spot weaknesses, fraud and waste.")],
 kicker="ACCOUNTING TYPES",cols=2,size=13)
img_slide("The Accounting Cycle","accounting-cycle.png",kicker="FROM TRANSACTION TO STATEMENTS",
 caption="Identify transactions → journal entries → ledger → trial balance → adjustments → financial statements → close.")
content("Accounting Periods",[
 "An accounting period is the time span for which financial statements are prepared — e.g. a calendar year, a quarter or a month.",
 "The income statement and cash flow statement report amounts occurring DURING the period.",
 "The balance sheet reports assets and liabilities AS AT the final moment of the period."],kicker="TIMEFRAMES",size=18)
img_text_slide("Chart of Accounts","chart-of-accounts.jpg",
 [("A chart of accounts lists the account names a company uses to record transactions in its general ledger.",0),
  ("Typical numbering: assets, liabilities, equity, revenue and expenses.",0),
  ("Each account maps to either the balance sheet or the income statement.",0),
  ("Companies tailor the chart to their needs and add accounts as required.",0)],
 kicker="THE LEDGER STRUCTURE",size=15)
two_col("Single-Entry vs Double-Entry Accounting",[
 ("Single entry",0),
 ("One entry per transaction, centred on the income statement",1),
 ("Tracks cash receipts and disbursements — the cash book",1),
 ("Assets and liabilities are not tracked systematically",1),
 ("Suitable only for the simplest businesses",1)],
 [("Double entry",0),
 ("Every transaction hits at least two accounts",1),
 ("Debits always equal credits",1),
 ("The accounting equation stays in balance",1),
 ("The standard for every modern accounting system",1)],
 kicker="RECORDING SYSTEMS",lhead="Single entry",rhead="Double entry")
img_slide("Double-Entry Rules — Debits & Credits","debit-credit-rules.jpg",kicker="WHICH SIDE INCREASES?",
 caption="Assets and expenses increase with debits; liabilities, equity and revenue increase with credits.")
img_slide("Double-Entry Bookkeeping — Journal Example","double-entry-journal.png",kicker="GENERAL JOURNAL",
 caption="A $10,000 cash sale: debit Cash (A/C 301), credit Sales (A/C 401).")
content("Bookkeeping",[
 "Bookkeeping is the systematic gathering and recording of a company's financial transactions.",
 "Transactions are identified, approved, sorted and stored so they can be retrieved and presented in reports.",
 "Examples: purchase of supplies with cash · purchase/sale of merchandise on credit · rent · salaries and wages · buying equipment · borrowing from a bank."],
 kicker="THE RECORDING DISCIPLINE",size=17)
cards3("The Three Key Financial Statements",[
 (BLUE,"Balance Sheet",["Statement of financial position","Assets = Liabilities + Equity","Liquidity and net worth","A snapshot AS AT a date"]),
 (TEAL,"Income Statement",["Profit & Loss statement","Profits = Revenues − Expenses","Profitability over a period","Top line to bottom line"]),
 (VIOLET,"Cash Flow Statement",["Movement of cash in the business","Operating · investing · financing","Where cash came from and went","Reconciles to the cash balance"])],kicker="FINANCIAL REPORTS")
tile_grid("Users of Financial Statements",[
 ("Management","Run the business — plans, budgets and performance reviews."),
 ("Investors & owners","Assess returns, growth and the safety of their capital."),
 ("Lenders & creditors","Judge the ability to repay loans and trade credit."),
 ("Regulators & tax authorities","Compliance — ACRA filings, IRAS tax computations."),
 ("Suppliers & customers","Decide whether to extend credit or rely on the firm."),
 ("Employees","Job security, bonuses and profit-sharing.")],
 kicker="WHO READS THE NUMBERS",cols=2,size=13)

# ---- Balance sheet
big_statement("Assets = Liabilities + Owners' Equity",
 "The balance sheet shows what the business owns (utilisation of funds) against what it owes and what the owners put in (sources of funds).",
 "BALANCE SHEET · THE FUNDAMENTAL EQUATION",color=TEAL)
img_slide("The Accounting Equation","accounting-equation.png",kicker="ASSETS = LIABILITIES + EQUITY",
 caption="Assets are what I own; liabilities are what I owe; equity is what is left for the owners.")
tile_grid("Balance Sheet — Definitions",[
 ("Assets","Resources that will provide a benefit in the future — utilisation of funds."),
 ("Liabilities","Obligations to repay money or provide a service in the future."),
 ("Owners' Equity","Money provided to the company by the owners."),
 ("Retained Earnings","Net income that is re-invested in the business.")],
 kicker="THE BUILDING BLOCKS",cols=2,size=14)
two_col("Balance Sheet Components",[
 ("Assets",0),
 ("Current assets: cash & equivalents, receivables, inventories, marketable securities",1),
 ("Long-term assets: plant, property & equipment (PP&E), intangible assets",1)],
 [("Liabilities & Equity",0),
 ("Current liabilities: accounts payable, current debt, current portion of LT debt",1),
 ("Long-term liabilities: bonds payable, long-term debt",1),
 ("Shareholders' equity: share capital + retained earnings",1)],
 kicker="STRUCTURE",lhead="Assets",rhead="Liabilities & Equity")
img_slide("Balance Sheet — Example","balance-sheet-example.png",kicker="XYZ COMPANY",
 caption="Current and long-term assets vs current and long-term liabilities and owners' equity.")
content("Balance Sheet — Limitations",[
 "The balance sheet is a static document — it depicts assets and liabilities 'as on date X' only.",
 "Assets are expressed at book value (historic cost less depreciation), which may differ from market value.",
 "Not all assets appear — e.g. internally-generated intangibles such as patents and a skilled workforce."],
 kicker="READ WITH CARE",size=18)
tile_grid("Current Assets",[
 ("Cash & equivalents","Legal tender, bills, coins, cheques; banker's acceptances, commercial paper, securities maturing < 90 days."),
 ("Accounts receivable","Funds owed by customers — recorded once goods or services are delivered."),
 ("Inventories","Raw materials, work-in-progress and finished goods held for sale."),
 ("Marketable securities","Common stock, treasury bills and bonds maturing within one year.")],
 kicker="LIQUID · QUICK TO CONVERT TO CASH",cols=2,size=13)
tile_grid("Long-Term Assets",[
 ("PP&E","Plant, property and equipment — the tangible productive base of the business."),
 ("Depreciation","Spreads the asset's cost over time to match the income it produces; the method chosen affects results."),
 ("Intangible assets","Patents, trademarks, goodwill — non-physical assets with economic value."),
 ("Working capital","Current assets − current liabilities: the funds available to run daily operations.")],
 kicker="NON-CURRENT ASSETS",cols=2,size=13)
formula_slide("Four Depreciation Methods",[
 ("Straight line","(Cost − Salvage) / Useful life","$50,000 over 10 years → $5,000/yr"),
 ("Double declining","2 × SL rate × Book value","Higher expense in early years"),
 ("Units of production","(Units / Life units) × (Cost − Salvage)","3M of 100M units on $50,000 → $1,500"),
 ("Sum of years","(Remaining life / SYD) × (Cost − Salvage)","$25,000, 8 yrs → yr 1 $5,556, yr 2 $4,861")],
 kicker="DEPRECIATION OF PP&E",accent=TEAL)
two_col("Liabilities",[
 ("Current liabilities",0),
 ("Accounts payable — money owed to vendors and suppliers",1),
 ("Current debt / notes payable — borrowings payable within one year",1),
 ("Current portion of long-term debt accrued this operating cycle",1)],
 [("Long-term liabilities",0),
 ("Bonds payable — principal at a future date, periodic interest",1),
 ("Bond types: serial, sinking fund, convertible, registered, secured, debenture",1),
 ("Long-term debt — owed to creditors beyond 12 months",1)],
 kicker="OBLIGATIONS",lhead="Current (< 1 year)",rhead="Long-term (> 1 year)")
tile_grid("Shareholders' Equity",[
 ("Share capital","Money invested by the owners — common or preferred shares."),
 ("Capital stages","Authorized → issued → subscribed → called-up → paid-up capital."),
 ("Retained earnings","Net profits kept in the company after all dividends and distributions."),
 ("Private equity","Ownership interest in a private company.")],
 kicker="THE OWNERS' STAKE",cols=2,size=14)

# ---- Income statement
big_statement("Profits = Revenues − Expenses",
 "The income statement (Profit & Loss) shows financial performance across an accounting period, from the top line to the bottom line.",
 "INCOME STATEMENT · THE FUNDAMENTAL EQUATION",color=VIOLET)
img_text_slide("Income Statement — Structure","income-statement-example.png",
 [("Revenue — the total of all sales (top line)",0,True),
  ("− Cost of Goods Sold: direct materials, labour, overheads",1),
  ("= Gross Profit",0,True),
  ("− Operating expenses: rental, admin payroll, utilities",1),
  ("= Operating Profit (EBIT)",0,True),
  ("− Interest expenses on debt",1),
  ("= Profit before Tax → − Tax = Profit after Tax",0,True)],
 kicker="FROM TOP LINE TO BOTTOM LINE",size=14)
tile_grid("Income Statement — Key Terms",[
 ("Revenue","Money actually received in the period, net of discounts and returns — the top line."),
 ("Expenses","Economic costs incurred to earn revenue; tax-deductible expenses lower taxable income."),
 ("EBITDA","Earnings before interest, tax, depreciation and amortization — operating performance without financing/accounting effects."),
 ("Depreciation & amortization","Wear-and-tear of tangible assets (depreciation) and write-down of intangibles (amortization)."),
 ("Tax","A mandatory financial charge imposed by the government on income."),
 ("Net income","What remains for the owners after every expense — the bottom line.")],
 kicker="COMPONENTS",cols=2,size=13)
flow_h("Summary of Income Statement Items",[
 "Revenue − COGS = Gross Profit",
 "− Marketing, selling & admin = EBITDA",
 "− Depreciation & amortization = EBIT",
 "− Interest & other expenses = EBT",
 "− Tax = Net Income"],kicker="THE PROFIT CASCADE",color=VIOLET)
img_slide("Cash vs Accrual Basis Accounting","accrual-vs-cash.png",kicker="WHEN IS REVENUE RECOGNIZED?",
 caption="Accrual basis: recognise when the work is done. Cash basis: recognise when the cash is received.")

# ---- Cash flow statement
tile_grid("Cash Flow Statement",[
 ("What it shows","The movement of cash in the business — incoming and outgoing money."),
 ("Operating activities","Converts income-statement items from accrual to cash: sales income and production expenses."),
 ("Investing activities","Purchase and sale of long-term investments and PP&E — gains or losses on assets."),
 ("Financing activities","Issuance/repurchase of stock and bonds, loans and dividend payments.")],
 kicker="FOLLOW THE CASH",cols=2,size=14)
img_slide("Cash Flow Statement — Example","cash-flow-statement-example.png",kicker="WORKED EXAMPLE",
 caption="Operating, investing and financing blocks reconcile the opening cash balance to the closing balance.")
A=ACTIVITIES[0]
mark("activity1")
activity_overview("ACTIVITY 1",A["title"],A["desc"],A["build"],A["services"],kicker="TOPIC 01 · HANDS-ON")
content(f"Recap — {T['title']}",[
 "The chart of accounts and double-entry rules keep the books in balance.",
 "The balance sheet shows position: Assets = Liabilities + Equity.",
 "The income statement shows performance: Profits = Revenues − Expenses.",
 "The cash flow statement shows where the cash actually moved.",
 "You prepared a full cash flow statement from a balance sheet and income statement."],kicker="TOPIC RECAP",size=17)

# ================================================== TOPIC 2
T=C.TOPICS[1]
mark("topic2")
section(f"TOPIC {T['code']}",T["title"],T["code"],T["subtitle"])
tile_grid(f"Key Concepts — {T['title']}",T["concepts"],kicker=f"TSC {T['weighting']}",cols=2,size=14)
content("What is Financial Ratio Analysis?",[
 "Financial ratios are mathematical indicators calculated by comparing key figures in the financial statements.",
 "They explain the reasons behind the current financial position and recent performance — and build expectations for the future.",
 "Example: net profit margin compares net income with net revenue — the dollars of profit earned per $100 of sales.",
 "Ratios only become meaningful against a benchmark: peers, the industry, or the company's own trend."],
 kicker="THE ANALYST'S TOOLKIT",size=17)
img_slide("Identifying Relevant Data","comparative-ratio-analysis.png",kicker="COMPARATIVE ANALYSIS",
 caption="Profitability, liquidity, leverage, efficiency and market-value ratios compared across companies.")
tile_grid("Benefits of Financial Ratios",[
 "Understand the profitability of the business.","Analyse operational efficiency.",
 "Identify the liquidity of the business.","Identify business risks.",
 "Identify financial risks.","Plan for the future."],
 kicker="WHY RATIOS",cols=2,size=14)
cards3("The Four Ratio Families",[
 (BLUE,"Liquidity & Leverage",["Current, quick, cash ratios","Operating cash flow ratio","Debt & debt-to-equity ratios","Interest & debt service coverage"]),
 (TEAL,"Efficiency",["Asset turnover","Inventory turnover","Receivables & payables turnover","Collection & payment periods"]),
 (VIOLET,"Profitability",["Gross / operating / net margins","Return on assets & equity","Return on capital employed","Earnings per share"])],kicker="RATIO MAP")

formula_slide("Liquidity Ratios",[
 ("Current Ratio","Current Assets / Current Liabilities","≥ 1 means current assets cover current bills"),
 ("Quick (Acid Test) Ratio","(Cash + Securities + Receivables) / CL","Most readily available assets only"),
 ("Cash Ratio","(Cash + Cash Equivalents) / CL","Normal value is below 1.00"),
 ("Operating Cash Flow Ratio","Operating Cash Flow / CL","Cash generated vs short-term obligations")],
 kicker="CAN THE FIRM PAY ITS BILLS?",accent=BLUE)
content("Reading the Liquidity Ratios",[
 "Current ratio ≥ 1: current assets exceed current liabilities — no expected liquidity problem; below 1 signals possible stress.",
 "Quick ratio strips out inventory; a rule of thumb is 0.5, but always compare with the industry.",
 "A quick ratio below the industry average may mean difficulty honouring current obligations.",
 "Creditors like a high cash ratio — but idle cash earns nothing, so firms typically hold less than 1.00."],
 kicker="INTERPRETATION",size=17)
img_slide("Quick Ratio — Worked Example","quick-ratio-example.png",kicker="COMPANY A VS COMPANY B",
 caption="Quick assets over current liabilities: Company A 0.95 vs Company B 1.21.")
formula_slide("Leverage Ratios",[
 ("Debt Ratio","Total Debt / Total Assets","< 0.5 stable · > 0.5 stability issue"),
 ("Debt-to-Equity","Total Liabilities / Shareholders' Equity","1.00 = half debt-financed; lower is safer"),
 ("Times Interest Earned","EBIT / Interest Expense","Can earnings cover the interest bill?"),
 ("Fixed Charge Coverage","(EBIT + Lease pmts) / (Interest + Lease pmts)","Adds lease obligations to the test"),
 ("Debt Service Coverage","Operating Income / Total Debt Service","Principal + interest coverage")],
 kicker="HOW MUCH RISK IN THE CAPITAL STRUCTURE?",accent=VIOLET)
content("Reading the Leverage Ratios",[
 "Debt ratio measures the risk that assets cannot pay off debts — critical for long-term sustainability.",
 "Very low debt can also mean under-utilised financing and restricted growth.",
 "A rising debt-to-equity trend is a warning: more of the assets are financed by lenders.",
 "Higher times-interest-earned is favourable — the company earns enough to pay its interest."],
 kicker="INTERPRETATION",size=17)
formula_slide("Efficiency Ratios (1)",[
 ("Asset Turnover","Net Sales / Average Total Assets","Revenue per dollar of assets"),
 ("Inventory Turnover","COGS / Average Inventories","High = efficient; low = over-stocking"),
 ("Receivables Turnover","Net Credit Sales / Average AR","How fast credit sales are collected"),
 ("Payables Turnover","Net Purchases / Average AP","How quickly suppliers are repaid")],
 kicker="HOW HARD ARE THE ASSETS WORKING?",accent=TEAL)
formula_slide("Efficiency Ratios (2)",[
 ("Fixed Asset Turnover","Net Revenue / Average Fixed Assets","Revenue per dollar of fixed assets"),
 ("Working Capital Turnover","Revenue / Average Working Capital","Revenue per dollar of working capital"),
 ("Average Collection Period","365 / Receivables Turnover","AR $40k on sales $400k → 36.5 days"),
 ("Average Payment Period","365 / Payables Turnover","AP turnover 2 → 182.5 days to pay")],
 kicker="TURNOVER & DAYS",accent=TEAL)
img_slide("Inventory Turnover — Worked Example","inventory-turnover-example.png",kicker="COMPANY A VS COMPANY B",
 caption="COGS over average inventory: Company A turns 2.33× vs Company B 1.75×.")
formula_slide("Profitability Ratios",[
 ("Gross Profit Margin","Gross Profit / Revenue","Production efficiency / pricing strategy"),
 ("Operating Margin","Operating Income (EBIT) / Revenue","9% → $0.09 profit per $1 of sales"),
 ("Net Profit Margin","Net Income / Net Sales","The most basic profitability measure"),
 ("Return on Assets","Net Income / Average Total Assets","Cents earned per dollar of assets"),
 ("Return on Capital Employed","Net Operating Profit / Capital Employed","Includes long-term finance"),
 ("Return on Equity","Net Income / Average Shareholders' Equity","The owners' return on new investment")],
 kicker="IS THE BUSINESS EARNING ENOUGH?",accent=AMBER)
content("Reading the Profitability Ratios",[
 "Higher ratios are more favourable — they reflect the ability to generate earnings.",
 "Compare only within the same industry: asset-heavy industries naturally show lower ROA.",
 "A rising ROA / ROE trend means profitability is improving; a falling trend means it is deteriorating.",
 "ROE can be manipulated with debt — never rely on ROE alone for investment decisions."],
 kicker="INTERPRETATION",size=17)
img_slide("Net Profit Ratio — Worked Example","net-profit-ratio-example.png",kicker="COMPANY A VS COMPANY B",
 caption="Net profit over net sales: Company A 26% vs Company B 25%.")
img_text_slide("Earnings Per Share (EPS)","eps-example.png",
 [("EPS = (Net Income − Preferred Dividends) / Weighted Average Common Shares",0,True),
  ("Standardises earnings by the number of shares outstanding.",0),
  ("Not comparable alone — different companies have different share counts.",0),
  ("Feeds the price-to-earnings (P/E) ratio used by investors.",0)],
 kicker="PER-SHARE PROFITABILITY",size=15)
img_text_slide("Measuring Profitability — Sample Company","sample-company-statements.png",
 [("Gross Margin: 37.50%",0,True),("Operating Margin: 26.67%",0,True),
  ("Net Profit Margin: 16.67%",0,True),("Return on Equity: 9.09%",0,True),
  ("Return on Assets: 4%",0,True),
  ("Derived from the sample balance sheet and income statement shown.",0)],
 kicker="PUTTING IT TOGETHER",size=15)
tile_grid("Statement of Changes in Equity",[
 ("What it is","Also called the statement of retained earnings — reconciles opening to closing equity."),
 ("The equation","Beginning Equity + Net Income − Dividends ± Other changes = Ending Equity."),
 ("Components","Opening balance · policy changes · prior-period corrections · restated balance."),
 ("More components","Changes in share capital · dividends · income/loss · revaluation reserve · other gains and losses.")],
 kicker="K4 · EQUITY CHANGES STATEMENT",cols=2,size=13)
two_col("Summary of All Ratios",[
 ("Liquidity",0),
 ("Current = CA / CL",1),("Acid test = (CA − Inventories) / CL",1),
 ("Cash = Cash & equivalents / CL",1),("Operating CF = OCF / CL",1),
 ("Leverage",0),
 ("Debt = Total liabilities / Total assets",1),
 ("Debt-to-equity = Total liabilities / Equity",1),
 ("Debt service coverage = Op. income / Debt service",1)],
 [("Efficiency",0),
 ("Asset turnover = Net sales / Avg total assets",1),
 ("Inventory turnover = COGS / Avg inventory",1),
 ("Receivables turnover = Net sales / Avg AR",1),
 ("Payables turnover = Net purchases / Avg AP",1),
 ("Profitability",0),
 ("Gross margin = Gross profit / Net sales",1),
 ("Operating margin = EBIT / Net sales · Net margin = NI / Net sales",1),
 ("ROE = NI / Equity · ROA = NI / Total assets",1)],
 kicker="ONE-PAGE CHEAT SHEET",lhead="Liquidity & Leverage",rhead="Efficiency & Profitability")
A=ACTIVITIES[1]
mark("activity2")
activity_overview("ACTIVITY 2",A["title"],A["desc"],A["build"],A["services"],kicker="TOPIC 02 · HANDS-ON")
A=ACTIVITIES[2]
mark("activity3")
activity_overview("ACTIVITY 3",A["title"],A["desc"],A["build"],A["services"],kicker="TOPIC 02 · HANDS-ON")
content(f"Recap — {T['title']}",[
 "Liquidity, leverage, efficiency and profitability — the four ratio families.",
 "Every ratio needs a benchmark: industry, competitor or the company's own trend.",
 "The statement of changes in equity reconciles opening to closing equity.",
 "You compared two companies ratio-by-ratio and ran a four-year trend analysis."],kicker="TOPIC RECAP",size=17)

# ================================================== TOPIC 3
T=C.TOPICS[2]
mark("topic3")
section(f"TOPIC {T['code']}",T["title"],T["code"],T["subtitle"])
tile_grid(f"Key Concepts — {T['title']}",T["concepts"],kicker=f"TSC {T['weighting']}",cols=2,size=14)
big_statement("Budgeting is balancing your expenses with your income —","by creating a plan to spend your money. To balance a budget you have two levers: increase income or decrease spending.","BUDGET",color=AMBER)
tile_grid("Key Components of a Budget",[
 ("Fixed expenses","Stay the same month to month — e.g. rent payments."),
 ("Flexible expenses","Change month to month — e.g. utilities."),
 ("Capital vs operating","Capital expenses buy long-term assets; operating expenses run the business."),
 ("Total income","Income from operations and other sources — investments, rentals and more."),
 ("Disposable income","What is left after taxes are subtracted."),
 ("Variance","The gap between the plan and the actual outcome.")],
 kicker="EXPENSES · INCOME",cols=2,size=13)
tile_grid("Broad Budgeting Methods",[
 ("Baseline budget","Begins with a previous plan as the baseline."),
 ("Incremental budget","A % or $ increment on the previous baseline."),
 ("Zero-based budget","Starts fresh, like a brand-new plan."),
 ("Hybrid budget","A combination of any of the above.")],
 kicker="FOUR APPROACHES",cols=2,size=14)
tile_grid("Budgeting Vocabulary",[
 "Fiscal year · budget communications","Budget targets · budget guidelines",
 "Annual plan · budget review","Board review · approved plan",
 "Forecast · actual","Variance"],
 kicker="TERMS YOU WILL MEET",cols=2,size=14)
content("Budgeting — Objectives",[
 "Provides clear structure, guidance and a basis for the action plan.",
 "Predicts expected cash flows and allocates resources to the right activities.",
 "Lets management prioritise projects and model scenarios — worst, median, best and custom cases.",
 "Measures actual performance against the approved plan and analyses the causes of variances.",
 "Ensures control over spending and mandates responsibility and accountability."],
 kicker="WHY BUDGET",size=16)
two_col("Budgeting — Advantages & Success Factors",[
 ("Advantages",0),
 ("Planning orientation and profitability review",1),
 ("Assumptions review and performance evaluation",1),
 ("Funding planning and cash allocation",1),
 ("Bottleneck analysis",1)],
 [("Success factors",0),
 ("Management support and employee involvement",1),
 ("Clear organisational goals and responsibility accounting",1),
 ("Sound accounting system and organisational structure",1),
 ("Flexibility and communication of results",1)],
 kicker="MAKING BUDGETS WORK",lhead="Advantages",rhead="Success factors")
img_slide("The Budgeting Process","budget-process.jpg",kicker="FROM STAKEHOLDER INPUT TO MONITORING",
 caption="Stakeholder input → goal setting → needs assessment → direction → operating budgets → adoption → monitoring → adjustment.")
img_slide("Sample Budget Report","sample-budget-report.png",kicker="BUDGET OVERVIEW · FY PLAN",
 caption="A monthly phased budget: revenue and expense lines with totals for the fiscal year.")
tile_grid("Capital Budgeting",[
 ("Purpose","Planning for long-term investments in projects and assets."),
 ("Cash, not profit","Focus on cash inflows and outflows, not accounting revenue and expenses."),
 ("Time value of money","Adjust future cash flows to today's dollars before comparing."),
 ("Decision methods","Payback · discounted payback · NPV · profitability index · IRR.")],
 kicker="INVESTING FOR THE LONG TERM",cols=2,size=14)
formula_slide("Time Value of Money",[
 ("Present Value","PV = FV / (1 + i)^n","$10,000 in 2 yrs @5% → $9,070.30 today"),
 ("Future Value","FV = PV × (1 + i)^n","$10,000 today @5% → $11,025 in 2 yrs"),
 ("Components","PV · FV · N periods · Interest i · Payment PMT","Opportunity cost: interest or inflation")],
 kicker="A DOLLAR TODAY BEATS A DOLLAR TOMORROW",accent=AMBER)
flow_h("Capital Budgeting Process",[
 "Identify long-term goals",
 "Identify investment proposals",
 "Estimate & analyse the cash flows",
 "Analyse financial feasibility",
 "Choose & implement projects",
 "Monitor the projects"],kicker="SIX STEPS",color=AMBER)
formula_slide("Capital Budgeting Methods",[
 ("Payback Period","Outlay / Annual cash inflow","$2,000 / $375 = 5.33 years"),
 ("Discounted Payback","Cumulative discounted CF ≥ outlay","Same project at 10% → ≈ 8 years"),
 ("Net Present Value","Σ PV(inflows) − Outlay","Accept if NPV > 0"),
 ("Profitability Index","PV(inflows) / PV(outflows)","Accept if PI > 1"),
 ("Internal Rate of Return","Rate where NPV = 0","Accept if IRR > threshold rate")],
 kicker="FIVE DECISION TOOLS",accent=AMBER)
img_slide("Discounted Payback — Worked Example","discounted-payback-example.png",kicker="DISCOUNTING THE CASH FLOWS",
 caption="Each future inflow is discounted before accumulating — payback arrives later than the simple method.")
img_text_slide("Profitability Index — Worked Example","profitability-index-example.png",
 [("Discount every inflow at the required return.",0),
  ("PV of inflows $2,304 vs PV of outflows $2,000.",0),
  ("PI = 2,304 / 2,000 = 1.15 → accept the project.",0,True),
  ("NPV = $304 — the same accept signal.",0)],
 kicker="PI > 1 → ACCEPT",size=15)
A=ACTIVITIES[3]
mark("activity4")
activity_overview("ACTIVITY 4",A["title"],A["desc"],A["build"],A["services"],kicker="TOPIC 03 · HANDS-ON")
A=ACTIVITIES[4]
mark("activity5")
activity_overview("ACTIVITY 5",A["title"],A["desc"],A["build"],A["services"],kicker="TOPIC 03 · HANDS-ON")

img_text_slide("Cash Flow Analysis & Management","monthly-cash-flow.png",
 [("Track the monthly statement of cash flows.",0),
  ("Improve cash flow: invoice promptly and chase receivables.",0),
  ("Negotiate supplier terms; manage inventory tightly.",0),
  ("Time capital spending; keep a cash buffer.",0)],
 kicker="MONTHLY CASH DISCIPLINE",size=15)
tile_grid("Pro Forma (Projected) Financial Statements",[
 ("Purpose","Plan ahead and compare actual vs predicted performance."),
 ("What is projected","The income statement and the balance sheet."),
 ("Cash flow","The projected cash flow is derived from the other two statements."),
 ("Forecasting cash","Estimate monthly sales → predict receipt of payments → estimate costs.")],
 kicker="FORECASTING THE STATEMENTS",cols=2,size=14)
content("Budget, Actual, Variance, Forecast",[
 "Budget: planned sales of $500,000.  Actual: sales of $400,000.",
 "Variance: −$100,000 — or −20% of budget ($100,000 / $500,000).",
 "Both the dollar and percentage variance appear on the Budget-to-Actual report.",
 "Forecast: the projection of future periods from historical data, updated at regular intervals."],
 kicker="THE CONTROL LOOP",size=17)
two_col("Budget vs Forecast",[
 ("Budget",0),
 ("A financial plan expressed in quantitative terms, prepared in advance",1),
 ("The financial expression of a business target",1),
 ("Sets targets · updated annually",1),
 ("Variance analysis: yes",1)],
 [("Forecast",0),
 ("Estimation of future trends from past and present data",1),
 ("A prediction of upcoming events and trends",1),
 ("No targets · updated at regular intervals",1),
 ("Variance analysis: no",1)],
 kicker="PLAN VS PREDICTION",lhead="Budget",rhead="Forecast")
img_text_slide("Types of Variances","variance-threshold.jpg",
 [("Adverse variance: actual income below budget, or spending above budget — a deficit.",0),
  ("Favourable variance: actual income above budget, or spending below budget — a surplus.",0),
  ("Manage over/underspend with clear-cut thresholds in dollars or percent.",0),
  ("RAG-status thresholds trigger review at defined limits.",0)],
 kicker="ADVERSE VS FAVOURABLE",size=15)
tile_grid("Business & Financial Risk",[
 ("High-risk company","Operates in a high-risk industry AND faces a risk of financial failure."),
 ("High-risk factors","Financial risk (banks, merchant accounts) · compliance · safety and statutory law."),
 ("Consequences","High-risk borrowers pay higher loan rates and larger down payments."),
 ("Risk types","Strategic · compliance · financial · operational · reputational risks."),
 ("Examples","Agriculture, construction and mining industries."),
 ("Low-risk company","Stable industry and sound finances — the reverse conditions.")],
 kicker="RISK PROFILING",cols=2,size=13)
cards3("Determine the Financial Health of a Company",[
 (BLUE,"Balance sheet",["Debt relative to equity","Short-term liquidity","Tangible vs financial assets","Collection & repayment cycles","Time to sell inventory"]),
 (TEAL,"Income statement",["Revenue growth by period","Gross profit margin","Net profit % of revenue","Interest cover on debt","Payout vs reinvestment"]),
 (VIOLET,"Cash flow statement",["Liquidity situation","Sources of cash","Free cash flow generated","Overall cash increase or decrease"])],kicker="THREE-STATEMENT REVIEW")
tile_grid("…then Ratio-Check the Health",[
 ("Margins","Gross profit margin and net profit margin on revenue."),
 ("Coverage","Ability to meet obligations — debt and interest payments."),
 ("Liquidity","Current ratio and quick ratio for obligations under one year."),
 ("Leverage","Debt-to-equity: the financing mix."),
 ("Efficiency","Inventory turnover and total asset turnover."),
 ("Returns","ROE and ROA: profit from equity and from assets.")],
 kicker="RATIO CHECKLIST",cols=2,size=13)
tile_grid("Financial Analysis for Investment Suitability",[
 ("Revenues","Quantity, quality and timing · growth (ignore one-offs) · concentration · revenue per employee."),
 ("Profits","Gross, operating and net profit margins."),
 ("Operational efficiency","Receivables turnover and inventory turnover."),
 ("Market data","P/E ratio · price-to-book · PEG ratio."),
 ("Capital & solvency","Return on equity · debt-to-equity · liquidity."),
 ("DuPont index","Net Profit Margin × Asset Turnover × Equity Multiplier = ROE decomposed.")],
 kicker="WHAT INVESTORS EXAMINE",cols=2,size=13)
two_col("Industry Ratios — Banking & Insurance",[
 ("Banking",0),
 ("Net interest margin = (Interest income − expense) / Total assets",1),
 ("Efficiency ratio = Non-interest expense / Revenue",1),
 ("Liquidity coverage ratio · leverage ratio · CET1 ratio",1)],
 [("Insurance",0),
 ("Persistency = paying policyholders / active policyholders",1),
 ("Solvency = available / required solvency margin",1),
 ("Combined ratio · incurred claims ratio · claim settlement ratio",1)],
 kicker="SECTOR-SPECIFIC RATIOS",lhead="Banking",rhead="Insurance")
img_slide("Benchmarking Financial Performance","benchmarking-sectors.png",kicker="AGAINST INDUSTRY AVERAGES",
 caption="Trend of US sectoral averages — compare your ratios against the sector, not in isolation.")
cards3("Financial Statement Analysis — Methods",[
 (BLUE,"Ratio analysis",["Generate the ratios","Compare against benchmarks","Track the ratio trend over time"]),
 (TEAL,"Horizontal (trend)",["Compare across periods","% change of the same line item","Spot growth and deterioration"]),
 (VIOLET,"Vertical analysis",["Every line as % of a base figure","IS lines as % of gross sales","BS lines as % of total assets"])],kicker="THREE LENSES")
img_slide("Horizontal Analysis — Example","horizontal-analysis-example.jpg",kicker="% CHANGE ACROSS PERIODS",
 caption="Income statement lines for two years with the percentage change of each item.")
img_slide("Horizontal Analysis — Case Study","horizontal-case-bs.png",kicker="FOUR-YEAR BALANCE SHEET",
 caption="Balance-sheet items tracked across 2018–2021 — the raw material for trend analysis.")
img_slide("Vertical Analysis — Example","vertical-analysis-example.png",kicker="EVERY LINE AS % OF THE BASE",
 caption="Income-statement items expressed as a percentage of sales for three years.")
A=ACTIVITIES[5]
mark("activity6")
activity_overview("ACTIVITY 6",A["title"],A["desc"],A["build"],A["services"],kicker="TOPIC 03 · HANDS-ON")
content(f"Recap — {T['title']}",[
 "Budgets plan the spending; forecasts predict the outcome; variances close the loop.",
 "Capital budgeting decisions: payback, discounted payback, NPV, PI and IRR.",
 "Financial health = balance sheet + income statement + cash flow, then ratio-check.",
 "Ratio, horizontal and vertical analysis — plus benchmarking — complete the toolkit.",
 "You evaluated an investment project and compared two companies' solvency risk."],kicker="TOPIC RECAP",size=17)

# ---------------- CLOSE ----------------
mark("wrapup")
section("WRAP-UP","Course Summary & Next Steps","")
tile_grid("What You Achieved",[
 ("Financial statements (LO1)","Read and prepared balance sheets, income statements and cash flow statements."),
 ("Ratio analysis (LO2)","Evaluated performance with liquidity, leverage, efficiency and profitability ratios and their trends."),
 ("Planning & budgeting (LO3)","Built budgets, evaluated investments (payback, NPV, PI) and assessed financial health and risk.")],
 kicker="LEARNING OUTCOMES",cols=1,size=15)
content("Recommended Courses",C.RECOMMENDED_COURSES,kicker="KEEP LEARNING",size=17)
content("Support",[
 "If you have any enquiries during and after the class, contact us:",
 "Email: enquiry@tertiaryinfotech.com",
 "Tel: +65 6100 0613",
 "Website: www.tertiarycourses.com.sg"],kicker="WE'RE HERE TO HELP",size=18)
content("Assessment",[
 C.ASSESSMENT["written"], C.ASSESSMENT["practical"],
 "Open book: slides, Learner Guide and approved materials only.",
 "Remember to take the Assessment digital attendance (TRAQOM).",
 "Submit your completed answers on the LMS at https://lms-tms.tertiaryinfotech.com/."],kicker="WRAP-UP")
flow_h("Assessment Flow",[
 "TRAQOM survey — scan the QR code on the LMS",
 "Assessment digital attendance — scan the SSG QR",
 "Sit WA (SAQ) then PP — open book",
 "Submit your answers on the LMS",
 "Sign the Assessment Summary Record"],kicker="ON ASSESSMENT DAY")
img_slide("Cert & TRAQOM Survey (Mandatory)","cert-traqom.png",kicker="https://lms-tms.tertiaryinfotech.com/",
 caption="Scan the QR codes to receive your certificate and complete the TRAQOM survey.")
content("Digital Attendance (Mandatory)",[
 "It is mandatory to take the AM, PM and Assessment digital attendance for WSQ-funded courses.",
 "The trainer or administrator will show you the digital attendance QR code generated from the SSG portal.",
 "Scan the QR code with your mobile phone camera and submit your attendance.",
 "A minimum of 75% attendance is required to be eligible for assessment and funding."],kicker="TRAQOM · SSG DIGITAL ATTENDANCE")
big_statement("Thank You!","You can now read, analyse and plan with financial statements — and put them to work for your SME.","SEE YOU IN THE NEXT COURSE",color=TEAL)

MARKS["total"]=len(prs.slides._sldIdLst)
OUT=os.path.join(REPO,"courseware",f"{C.SHORT_TITLE}-{C.VERSION}.pptx")
prs.save(OUT)
import json
with open(os.path.join(HERE,"slide_index.json"),"w") as f:
    json.dump(MARKS,f,indent=1)
print(f"Saved {OUT}  ({MARKS['total']} slides)")
print("Slide index:",MARKS)
